import os
import yaml

from docx import Document
from pptx import Presentation
from pptx.util import Inches


doc_names_list = ['Документ на подпись по продажам', 'Страхование сотрудников', 'Список контрагентов', 'Документация на проект', 'НПА по отделам', 'Структура организации']
ppt_names_list = ['Презентация нового продукта 2023', 'Методические рекомнедации', 'Презентация продукции', 'Презентация внутренних ресурсов']
xls_name_list = ['Финансовый отчет за ноябрь', 'Список сотрудников', 'Положение о структурном подразделении', 'Блан протокола', 'Докладная записка', 'Приказ на создание ', 'Табель', 'Отчет по командировке']

win_path = ['']

class File_creator:

    #TODO WRITE FUNCTION WICH GENERATE FILES CONTENT
    def generate_text_content():
        pass

    #TODO WRITE PDF CREATION  
    def create_pdf(name, content=""):
        pass

    def create_powpoint(name, content=""):
        presentation = Presentation()

        slide_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Привет, мир!"
        subtitle.text = "Пример создания презентации PowerPoint с помощью python-pptx"

        slide_layout = presentation.slide_layouts[5]
        slide = presentation.slides.add_slide(slide_layout)
        img_path = 'img1.jpg'
        left = top = Inches(0)
        pic = slide.shapes.add_picture(img_path, left, top)

        presentation.save(name + '.pptx')
        print(f"Презентация {name + '.pptx'} - успешно создана.")

    def create_text(name, content=""):
        file = open(name + ".txt", "w+t")
        file.wtite(content)
        file.close

    def create_word(name, content=""):
        doc = Document()
        doc.add_heading(name + '.docx', level=1)
        doc.add_paragraph('Простой документ .docx из python.')

        table = doc.add_table(rows=3, cols=3)
        for i in range(3):
            for j in range(3):
                table.cell(i, j).text = f'Ячейка {i+1}-{j+1}'

        doc.save(name + '.docx')

        print(f"Документ {name + '.docx'} - успешно создан.")


