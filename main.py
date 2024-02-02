from docx import Document
from pptx import Presentation
from pptx.util import Inches


doc_names_list = ['Документ на подпись по продажам']
ppt_names_list = ['Презентация нового продукта 2023']
xls_name_list = ['Финансовый отчет за ноябрь']

#TODO WRITE FUNCTION WICH GENERATE FILES CONTENT
def generate_text_content():
    pass

#TODO WRITE FUNCTION WICH PLACES FILES AT DIRECTORIES USING SSH/
def place_files_remotly():
    pass

#TODO WRITE FUNCTION WICH PLACES FILES AT DIRECTORIES USING LOCAL
def place_files_localy():
    pass

def generate_path(file_type, id):
    path_list = ['']

    if file_type == 'xls':
        path = path_list[1][id]
    elif file_type == 'docx':
        path = path_list[2][id]
    elif file_type == 'pptx':
        path = path_list[3][id]
    elif file_type == 'txt':
        path = path_list[4][id]
    elif file_type == 'txt':
        path = path_list[5][id]
    return path

#TODO WRITE PDF CREATION  
def create_pdf(name, content=""):
    pass

def create_powpoint(name, content=""):
    presentation = Presentation()

    slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Привет, мир!"
    subtitle.text = "Пример создания презентации PowerPoint с помощью python-pptx"

    slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(slide_layout)
    img_path = 'img1.jpg'
    left = top = Inches(0)
    pic = slide.shapes.add_picture(img_path, left, top)

    presentation.save(name + '.pptx')
    print(f"Презентация {name + '.pptx'} - успешно создана.")

def create_text(name, content=""):
    file = open(name + ".txt", "w+t")
    file.wtite(content)
    file.close

def create_word(name, content=""):
    doc = Document()
    doc.add_heading(name + '.docx', level=1)
    doc.add_paragraph('Простой документ .docx из python.')

    table = doc.add_table(rows=3, cols=3)
    for i in range(3):
        for j in range(3):
            table.cell(i, j).text = f'Ячейка {i+1}-{j+1}'

    doc.save(name + '.docx')

    print(f"Документ {name + '.docx'} - успешно создан.")




for a in ppt_names_list:
    create_powpoint(a)