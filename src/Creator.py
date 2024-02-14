import os
import yaml

from docx import Document
from pptx import Presentation
from pptx.util import Inches


class File_creator:
    def __init__(self):
        pass

    def check_bd(self):
        pass

    def get_bd(self):
        print(self.parsed_names)
        print(self.parsed_paths)
 
    def parse_bd(self, type):
        self.parsed_names = []
        self.parsed_paths = []
        file = open("PayloadFiles/names.yaml", 'r')
        data = yaml.safe_load(file)
        for line, info in data.items():
            if line == type:
                for line2, data in info.items():
                    if line2 == 'names':
                        for ptr in range(len(data)):
                            print(data[ptr])
                            self.parsed_names.append(data[ptr])
                    elif line2 == 'paths':
                        for ptr in range(len(data)):
                            print(data[ptr])
                            self.parsed_paths.append(data[ptr])
        
    #TODO WRITE FUNCTION WICH GENERATE FILES CONTENT
    def generate_text_content(self, name):
        pass

    #TODO WRITE PDF CREATION  
    def create_pdf(self, name, content=""):
        pass

    def create_powpoint(self, name, content=""):
        presentation = Presentation()

        slide_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Привет, мир!"
        subtitle.text = "Пример создания презентации PowerPoint с помощью python-pptx"

        slide_layout = presentation.slide_layouts[5]
        slide = presentation.slides.add_slide(slide_layout)
        img_path = 'img1.jpg'
        left = top = Inches(0)
        pic = slide.shapes.add_picture(img_path, left, top)

        presentation.save(name + '.pptx')
        print(f"Презентация {name + '.pptx'} - успешно создана.")

    def create_text(self, name, content=""):
        file = open(name + ".txt", "w+t")
        file.wtite(content)
        file.close

    def create_word(self, name, content=""):
        doc = Document()
        doc.add_heading(name + '.docx', level=1)
        doc.add_paragraph('Простой документ .docx из python.')

        table = doc.add_table(rows=3, cols=3)
        for i in range(3):
            for j in range(3):
                table.cell(i, j).text = f'Ячейка {i+1}-{j+1}'

        doc.save(name + '.docx')

        print(f"Документ {name + '.docx'} - успешно создан.")


